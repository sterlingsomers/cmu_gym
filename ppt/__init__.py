from pptx import Presentation as ppt
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO_SHAPE

import os


DEFAULT_LEFT = Inches(0.5)
DEFAULT_TOP = Inches(1.75)
DEFAULT_HEIGHT = Inches(5.5)
DEFAULT_TABLE_LEFT = Inches(3.0)
DEFAULT_TABLE_TOP_TITLE_NEXT = Inches(1.75)
DEFAULT_TABLE_TOP_NEXT = Inches(2.15)
DEFAULT_TABLE_TOP_TITLE_LAST = Inches(4.4)
DEFAULT_TABLE_TOP_LAST = Inches(4.8)
DEFAULT_TABLE_HEIGHT = Inches(2.25)
DEFAULT_TABLE_WIDTH = Inches(7.0)
DEFAULT_FONT = Pt(14)
DEFAULT_TITLE_COLOR = RGBColor(0x00,0x00,0x00)
DEFAULT_FILL_COLOR = RGBColor(0xFF,0xFF,0xFF)
DEFAULT_TITLE_HEIGHT = Inches(1.0)
DEFAULT_CHART_TOP = DEFAULT_TABLE_TOP_LAST
DEFAULT_CHART_LEFT = DEFAULT_TABLE_LEFT
DEFAULT_CHART_WIDTH = DEFAULT_TABLE_WIDTH
DEFAULT_CHART_HEIGHT = DEFAULT_TABLE_HEIGHT
DEFAULT_OVERLAY_WIDTH = DEFAULT_TABLE_WIDTH

DEFAULT_TITLE_INFO_TOP = DEFAULT_TOP
DEFAULT_TITLE_INFO_LEFT = DEFAULT_LEFT
DEFAULT_TITLE_INFO_HEIGHT = DEFAULT_HEIGHT
DEFAULT_TITLE_INFO_WIDTH = DEFAULT_TABLE_WIDTH

def add_barchart_to_slide(
        title_top,
        slide,
        title,
        chart_dictionary,
        chart_order=[],
        chart_top=DEFAULT_CHART_TOP,
        chart_left=DEFAULT_CHART_LEFT,
        chart_width=DEFAULT_CHART_WIDTH,
        chart_height=DEFAULT_CHART_HEIGHT):
    chart_x = chart_left
    chart_y = chart_top
    chart_cx = chart_width
    chart_cy = chart_height
    title_shape = slide.shapes.add_textbox(
        left=chart_left,
        top=title_top,
        width=chart_width,
        height=DEFAULT_TABLE_HEIGHT
    )
    title_shape.text = title
    title_paragraph = title_shape.text_frame.paragraphs[0]
    title_paragraph.font.color.rgb = DEFAULT_TITLE_COLOR
    if chart_dictionary is None or not isinstance(chart_dictionary,dict):
        return
    if len(chart_order) == 0:
        chart_order = list(chart_dictionary.keys())

    # define chart data ---------------------
    chart_data = CategoryChartData()
    chart_data.categories = chart_order
    chart_values = []
    for i in range(len(chart_order)):
        chart_values.append(chart_dictionary[chart_order[i]])
    chart_data.add_series('Series 1', chart_values)

    # add chart to slide --------------------
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, chart_x, chart_y, chart_cx, chart_cy, chart_data
    )
    this_chart = chart_shape.chart
    category_axis = this_chart.category_axis
    category_axis.tick_labels.font.size = Pt(16)
    value_axis = this_chart.value_axis
    value_axis.tick_labels.font.color.rgb = DEFAULT_FILL_COLOR


def add_image_to_slide(title_top, image_top, slide, title, image_path):
    try:
        title_shape = slide.shapes.add_textbox(
            left=DEFAULT_TABLE_LEFT,
            top=title_top,
            width=DEFAULT_TABLE_WIDTH,
            height=DEFAULT_TABLE_HEIGHT
        )
        title_shape.text = title
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.font.color.rgb = DEFAULT_TITLE_COLOR

        slide.shapes.add_picture(image_path, left=DEFAULT_TABLE_LEFT, top=image_top, height=DEFAULT_TABLE_HEIGHT)


    except Exception as e:
        print("WARNING: bad table creation: " + str(e))

def add_table_to_slide(title_top, table_top, slide, title, table_matrix):
    try:
        title_shape = slide.shapes.add_textbox(
            left=DEFAULT_TABLE_LEFT,
            top=title_top,
            width=DEFAULT_TABLE_WIDTH,
            height=DEFAULT_TABLE_HEIGHT
        )
        title_shape.text = title
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.font.color.rgb = DEFAULT_TITLE_COLOR
        if table_matrix is None:
            return
        table_shape = slide.shapes.add_table(
            rows=5,
            cols=6,
            left=DEFAULT_TABLE_LEFT,
            top=table_top,
            width=DEFAULT_TABLE_WIDTH,
            height=DEFAULT_TABLE_HEIGHT
        )
        this_table = table_shape.table
        for row in range(len(table_matrix)):
            for col in range(len(table_matrix[row])):
                cell = this_table.cell(row, col)
                cell.text = table_matrix[row][col]
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.size = DEFAULT_FONT
                paragraph.alignment = PP_ALIGN.CENTER


    except Exception as e:
        print("WARNING: bad table creation: " + str(e))

def image_title(t, info):
    title = 'Step ' + str(t)
    if 'competency' in info:
        title += ': ' + info['competency']
    if 'status' in info:
        title += ': ' + info['status']
    if 'package_state' in info:
        title += ': ' + info['package_state']
    return title


class Presentation:
    def __init__(self, savepath, name):
        self.savepath = savepath
        self.name = name
        self.savename = os.path.join(savepath, name)
        self.presentation = ppt()

    def add_image_slide(
            self,
            image_path,
            title='',
            left=DEFAULT_LEFT,
            top=DEFAULT_TOP,
            height=DEFAULT_HEIGHT):
        slide_layout = self.presentation.slide_layouts[5]
        slide = self.presentation.slides.add_slide(slide_layout)
        slide.shapes.add_picture(image_path, left=left, top=top, height=height)
        slide_title = slide.shapes.title
        slide_title.text = title
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left=DEFAULT_TABLE_LEFT, top=top, width=DEFAULT_OVERLAY_WIDTH, height=height)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = DEFAULT_FILL_COLOR
        return slide

    def add_title_slide(
            self,
            title,
            hikerx,
            hikery,
            dronex,
            droney,
            map):
        # Title the Presentation
        slide_layout = self.presentation.slide_layouts[5]
        slide = self.presentation.slides.add_slide(slide_layout)
        slide_title = slide.shapes.title
        slide_title.text = title
        # Show basic mission data
        info_shape = slide.shapes.add_textbox(
            left=DEFAULT_TITLE_INFO_LEFT,
            top=DEFAULT_TITLE_INFO_TOP,
            width=DEFAULT_TITLE_INFO_WIDTH,
            height=DEFAULT_TITLE_INFO_HEIGHT
        )
        info =  'Map:      {}\n'.format(map)
        info += 'Hiker at: ({},{})\n'.format(hikerx,hikery)
        info += 'Drone at: ({},{})\n'.format(dronex,droney)
        info_shape.text = info
        title_paragraph = info_shape.text_frame.paragraphs[0]
        title_paragraph.font.color.rgb = DEFAULT_TITLE_COLOR

        return slide

    def save(self):
        self.presentation.save(self.savename)

